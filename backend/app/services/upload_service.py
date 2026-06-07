"""
Upload Service — orchestrates multimodal document ingestion pipeline
"""
from typing import Any
from datetime import datetime
from app.utils.logger import logger


class UploadService:
    def __init__(self, db=None):
        self.db = db

    async def process_document(self, doc_id: str, file_path: str, file_ext: str):
        """
        Full ingestion pipeline:
        1. Parse document by type
        2. Chunk content
        3. Generate embeddings
        4. Upsert to vector store
        5. Update DB status
        """
        try:
            logger.info(f"[UploadService] Processing {doc_id} ({file_ext})")
            await self._update_status(doc_id, "processing")

            # Step 1: Parse
            raw_chunks = await self._parse(file_path, file_ext)

            # Step 2: Chunk text portions
            from app.rag.chunking import RecursiveTextSplitter
            splitter = RecursiveTextSplitter(chunk_size=512, chunk_overlap=64)

            all_chunks = []
            for raw in raw_chunks:
                if raw.get("modality") == "text":
                    for chunk in splitter.split(raw["text"]):
                        all_chunks.append({
                            **chunk.to_dict(),
                            "document_id": doc_id,
                            "modality": "text",
                        })
                else:
                    # Images, tables, audio go in as-is
                    all_chunks.append({**raw, "document_id": doc_id})

            # Step 3: Save chunks to MongoDB
            if all_chunks and self.db:
                for i, chunk in enumerate(all_chunks):
                    chunk["_id"] = f"{doc_id}_{i}"
                    chunk["chunk_index"] = i
                    chunk["created_at"] = datetime.utcnow()
                await self.db["chunks"].insert_many(all_chunks, ordered=False)

            # Step 4: Embed text chunks and upsert to Qdrant
            text_chunks = [c for c in all_chunks if c.get("modality") == "text"]
            if text_chunks:
                try:
                    from app.rag.retriever import VectorRetriever
                    retriever = VectorRetriever()
                    await retriever.upsert(text_chunks, doc_id)
                except Exception as e:
                    logger.warning(f"[UploadService] Vector upsert failed: {e}")

            # Step 5: Update status
            await self._update_status(doc_id, "ready", chunk_count=len(all_chunks))
            logger.info(f"[UploadService] ✅ {doc_id} ingested: {len(all_chunks)} chunks")

        except Exception as e:
            logger.error(f"[UploadService] ❌ Failed {doc_id}: {e}")
            await self._update_status(doc_id, "failed")

    async def _parse(self, file_path: str, ext: str) -> list:
        """Route to appropriate parser based on file type."""
        if ext == "pdf":
            from app.multimodal.pdf_parser import PDFParser
            result = await PDFParser().parse(file_path)
            chunks = result["text_chunks"] + result["images"]
            for t in result["tables"]:
                text = "\n".join(" | ".join(str(c) for c in row) for row in t["data"] if row)
                chunks.append({"text": text, "modality": "table", "page": t["page"]})
            return chunks

        elif ext in ("png", "jpg", "jpeg", "gif", "webp"):
            from app.multimodal.image_processor import ImageProcessor
            img = await ImageProcessor().process(file_path)
            from app.multimodal.ocr_engine import OCREngine
            text = await OCREngine().extract_text(
                image_path=file_path, image_base64=img.get("image_base64"), mime_type=img.get("mime_type")
            )
            return [{"text": text or "Image content", "modality": "image", **img}]

        elif ext in ("mp3", "wav", "m4a"):
            from app.multimodal.audio_processor import AudioProcessor
            result = await AudioProcessor().transcribe(file_path)
            return [{"text": result.get("text", ""), "modality": "audio"}]

        elif ext in ("mp4", "mov", "avi"):
            from app.multimodal.video_processor import VideoProcessor
            result = await VideoProcessor().process(file_path)
            chunks = [{"text": result.get("transcript", ""), "modality": "audio"}]
            chunks += result.get("frames", [])
            return chunks

        elif ext in ("xlsx", "xls"):
            from app.multimodal.table_extractor import TableExtractor
            tables = await TableExtractor().extract_from_excel(file_path)
            return [{"text": t["text"], "modality": "table"} for t in tables]

        elif ext == "csv":
            from app.multimodal.table_extractor import TableExtractor
            tables = await TableExtractor().extract_from_csv(file_path)
            return [{"text": t["text"], "modality": "table"} for t in tables]

        elif ext in ("pptx", "ppt"):
            return await self._parse_pptx(file_path)

        elif ext in ("docx", "doc"):
            return await self._parse_docx(file_path)

        elif ext in ("txt", "md"):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            return [{"text": text, "modality": "text"}]

        return []

    async def _parse_pptx(self, file_path: str) -> list:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            chunks = []
            for slide_num, slide in enumerate(prs.slides):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                if texts:
                    chunks.append({"text": f"Slide {slide_num+1}:\n" + "\n".join(texts), "modality": "text"})
            return chunks
        except Exception as e:
            logger.error(f"[UploadService] PPTX parse failed: {e}")
            return []

    async def _parse_docx(self, file_path: str) -> list:
        try:
            from docx import Document
            doc = Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text)
            return [{"text": text, "modality": "text"}]
        except Exception as e:
            logger.error(f"[UploadService] DOCX parse failed: {e}")
            return []

    async def _update_status(self, doc_id: str, status: str, chunk_count: int = 0):
        if self.db:
            await self.db["documents"].update_one(
                {"_id": doc_id},
                {"$set": {"status": status, "chunk_count": chunk_count, "updated_at": datetime.utcnow()}},
            )
